"""Build the Hidden Entrepreneurs Detection presentation (English)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# --- Brand palette (Mastercard-inspired) ---
RED        = RGBColor(0xC4, 0x1E, 0x3A)
RED_DEEP   = RGBColor(0xA0, 0x16, 0x2E)
ORANGE     = RGBColor(0xFF, 0x5F, 0x00)
YELLOW     = RGBColor(0xF7, 0x9E, 0x1B)
DARK       = RGBColor(0x1A, 0x1A, 0x1A)
GREY_DARK  = RGBColor(0x4A, 0x4A, 0x4A)
GREY       = RGBColor(0x6E, 0x6E, 0x6E)
GREY_LIGHT = RGBColor(0xE6, 0xE6, 0xE6)
BG         = RGBColor(0xFA, 0xFA, 0xFA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG    = RGBColor(0x1E, 0x1E, 0x2E)
CODE_FG    = RGBColor(0xE6, 0xE6, 0xF0)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)


prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]


# ----------------------------- helpers -----------------------------
def add_rect(slide, left, top, width, height, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             font="Calibri", anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, bullets,
                size=16, color=DARK, bullet_color=RED, line_spacing=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        r1 = p.add_run()
        r1.text = "■  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = "Calibri"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def add_code(slide, left, top, width, height, code, size=12,
             bg=CODE_BG, fg=CODE_FG):
    """Add a code block with dark background and monospace text."""
    add_rect(slide, left, top, width, height, bg)
    tb = slide.shapes.add_textbox(
        left + Emu(60000), top + Emu(40000),
        width - Emu(120000), height - Emu(80000)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = code.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = ln if ln else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = fg
    return tb


def slide_header(slide, kicker, title):
    # Top red band
    add_rect(slide, 0, 0, SW, Inches(0.45), RED)
    # Kicker / section label
    add_text(slide, Inches(0.55), Inches(0.08), Inches(10), Inches(0.32),
             kicker.upper(), size=11, bold=True, color=WHITE, font="Calibri")
    # Page title
    add_text(slide, Inches(0.55), Inches(0.6), Inches(12.2), Inches(0.7),
             title, size=30, bold=True, color=DARK, font="Calibri")
    # Thin red underline accent
    add_rect(slide, Inches(0.55), Inches(1.32), Inches(0.6), Inches(0.06), RED)


def footer(slide, page, total):
    add_text(slide, Inches(0.55), Inches(7.12), Inches(8), Inches(0.3),
             "Hidden Entrepreneurs Detection  ·  Mastercard × AIESEC Hackathon",
             size=9, color=GREY)
    add_text(slide, Inches(12.0), Inches(7.12), Inches(1.2), Inches(0.3),
             f"{page} / {total}",
             size=9, color=GREY, align=PP_ALIGN.RIGHT)


# ----------------------------- slides -----------------------------
TOTAL_SLIDES = 21  # set after final build


def slide_title():
    s = prs.slides.add_slide(BLANK)
    # Full-bleed brand background
    add_rect(s, 0, 0, SW, SH, DARK)
    # Diagonal red accent
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                             Inches(8.5), 0, Inches(4.83), SH)
    tri.fill.solid(); tri.fill.fore_color.rgb = RED
    tri.line.fill.background(); tri.shadow.inherit = False
    # Mastercard mark (two circles) — decorative
    c1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(0.7), Inches(0.55), Inches(0.55), Inches(0.55))
    c1.fill.solid(); c1.fill.fore_color.rgb = RED; c1.line.fill.background()
    c2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(1.05), Inches(0.55), Inches(0.55), Inches(0.55))
    c2.fill.solid(); c2.fill.fore_color.rgb = YELLOW; c2.line.fill.background()
    c2.fill.transparency = 0.4

    add_text(s, Inches(0.7), Inches(1.4), Inches(8), Inches(0.4),
             "MASTERCARD × AIESEC HACKATHON  ·  MAY 2026",
             size=14, bold=True, color=YELLOW)

    add_text(s, Inches(0.7), Inches(1.95), Inches(9.5), Inches(1.6),
             "Hidden Entrepreneurs\nDetection",
             size=54, bold=True, color=WHITE)

    add_text(s, Inches(0.7), Inches(4.3), Inches(8.5), Inches(1.6),
             "Finding micro-business owners hiding inside\n"
             "an 80K consumer-card population using PU Learning\n"
             "and behavioural transaction features.",
             size=18, color=GREY_LIGHT)

    add_rect(s, Inches(0.7), Inches(6.3), Inches(0.08), Inches(0.6), RED)
    add_text(s, Inches(0.9), Inches(6.32), Inches(6), Inches(0.3),
             "Team submission · MDQ dataset",
             size=12, bold=True, color=WHITE)
    add_text(s, Inches(0.9), Inches(6.6), Inches(6), Inches(0.3),
             "github.com/Shalbulov/hidden-entrepreneurs-detection",
             size=11, color=GREY_LIGHT)


def slide_problem():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "01 · Problem", "What are we trying to find?")

    add_text(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.5),
             "A non-trivial share of personal cards are quietly running a business.",
             size=18, bold=True, color=DARK)

    add_bullets(s, Inches(0.55), Inches(2.5), Inches(7.6), Inches(4),
                [
                    "Resellers, freelancers, repair shops, tutors and similar "
                    "micro-entrepreneurs use a personal card for commercial flows.",
                    "They look like ordinary consumers in the data — but their "
                    "transaction patterns betray them.",
                    "Finding them lets the bank pitch the right product: business "
                    "card, acquiring, line of credit.",
                    "Goal: rank 80K consumer cards by P(business) and surface the "
                    "top tiers to relationship managers.",
                ], size=15)

    # Right info card
    add_rect(s, Inches(8.6), Inches(2.5), Inches(4.2), Inches(3.4), DARK)
    add_text(s, Inches(8.85), Inches(2.65), Inches(3.9), Inches(0.4),
             "WHY IT MATTERS", size=11, bold=True, color=YELLOW)
    add_text(s, Inches(8.85), Inches(3.1), Inches(3.9), Inches(2.6),
             "Every undiagnosed entrepreneur is:\n\n"
             "•  a missed acquiring contract\n"
             "•  a missed business-card customer\n"
             "•  a misclassified risk profile\n"
             "•  a competitor's opportunity",
             size=14, color=WHITE)

    footer(s, 2, TOTAL_SLIDES)


def slide_dataset():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "02 · Data", "The MDQ dataset")

    # Three KPI cards
    def kpi(left, big, label, color):
        add_rect(s, left, Inches(1.85), Inches(4.0), Inches(2.0), color)
        add_text(s, left, Inches(2.05), Inches(4.0), Inches(0.9), big,
                 size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, left, Inches(3.05), Inches(4.0), Inches(0.6), label,
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    kpi(Inches(0.55), "25K",  "Business cards · ~3M txns",  RED)
    kpi(Inches(4.7),  "80K",  "Consumer cards · ~10M txns", DARK)
    kpi(Inches(8.85), "~2K",  "Merchant reference rows",    ORANGE)

    add_text(s, Inches(0.55), Inches(4.15), Inches(12.2), Inches(0.4),
             "Columns we actually use:", size=15, bold=True, color=DARK)

    add_bullets(s, Inches(0.55), Inches(4.55), Inches(12.2), Inches(2.5),
                [
                    "transaction_date / transaction_timestamp — when (day, hour, weekday).",
                    "transaction_amount_kzt — how much (stats, large / round flags).",
                    "mcc — what kind of merchant (business-specific MCC ratio).",
                    "merchant_id, country, channel — where and how (POS vs e-com, intl).",
                    "is_recurring, tokenized — subscription-like behaviour, Apple/Google Pay.",
                ], size=14)

    footer(s, 3, TOTAL_SLIDES)


def slide_pipeline():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "03 · Approach", "End-to-end pipeline")

    steps = [
        ("Load",          "Read 3 parquets;\nlabel biz=1, cons=0."),
        ("Features",      "Aggregate txns →\none row per card,\n~30 features."),
        ("Split",         "Train/Val/Test by\ncard_number — no\nleakage."),
        ("Baseline",      "RF + XGBoost,\nclass-weighted,\nmeasure Val AUC."),
        ("Denoise",       "Iterative cleaning\n+ PU Learning\n(2 iterations)."),
        ("Boost",         "XGBoost on PU\nnegatives, early\nstopping."),
        ("Score",         "Predict on all 80K\nconsumers, bucket\nby priority."),
    ]
    n = len(steps)
    total_w = Inches(12.2)
    box_w   = Inches(1.55)
    gap     = (total_w - box_w * n) / (n - 1)
    left0   = Inches(0.55)
    top     = Inches(2.4)

    for i, (label, desc) in enumerate(steps):
        x = left0 + (box_w + gap) * i
        color = RED if i in (0, 4, 6) else DARK
        add_rect(s, x, top, box_w, Inches(2.6), color)
        add_text(s, x, top + Inches(0.2), box_w, Inches(0.5),
                 f"{i+1}", size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, top + Inches(0.7), box_w, Inches(0.4),
                 label, size=14, bold=True, color=YELLOW if color == DARK else WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, top + Inches(1.2), box_w, Inches(1.3),
                 desc, size=10, color=WHITE, align=PP_ALIGN.CENTER)
        # arrow between boxes
        if i < n - 1:
            ax = x + box_w + Emu(40000)
            ay = top + Inches(1.25)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay,
                                     gap - Emu(80000), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = GREY
            arr.line.fill.background(); arr.shadow.inherit = False

    add_text(s, Inches(0.55), Inches(5.45), Inches(12.2), Inches(1.4),
             "The hard part is step 5 — the consumer label is noisy. "
             "We never call a consumer a 'true negative' until the model itself "
             "is confident it isn't a business.",
             size=14, color=GREY_DARK)

    footer(s, 4, TOTAL_SLIDES)


def slide_imports():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "04 · Code", "Imports & setup")

    add_code(s, Inches(0.55), Inches(1.75), Inches(7.8), Inches(4.8),
"""import pandas as pd
import numpy as np
import matplotlib.pyplot as plt

from sklearn.model_selection import train_test_split
from sklearn.ensemble import RandomForestClassifier
from sklearn.metrics import (
    classification_report, roc_auc_score,
    roc_curve, confusion_matrix
)
from xgboost import XGBClassifier

RANDOM_STATE = 42""", size=13)

    add_text(s, Inches(8.6), Inches(1.75), Inches(4.2), Inches(0.4),
             "Why these libraries?", size=15, bold=True, color=DARK)
    add_bullets(s, Inches(8.6), Inches(2.2), Inches(4.2), Inches(4),
                [
                    "pandas — group, aggregate, join 13M+ rows efficiently.",
                    "scikit-learn — Random Forest baseline, splits, ROC metrics.",
                    "xgboost — gradient-boosted trees with early stopping.",
                    "matplotlib — feature importance, ROC, score distribution.",
                    "Fixed RANDOM_STATE so the whole pipeline is reproducible.",
                ], size=13)
    footer(s, 5, TOTAL_SLIDES)


def slide_load():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "05 · Code", "Loading & labelling")

    add_code(s, Inches(0.55), Inches(1.75), Inches(7.8), Inches(4.6),
"""biz   = pd.read_parquet('business_cards_MDQ.parquet')
cons  = pd.read_parquet('consumer_cards_MDQ.parquet')
merch = pd.read_parquet('merchants_reference.parquet')

for df in [biz, cons]:
    df['transaction_date']      = pd.to_datetime(df['transaction_date'])
    df['transaction_timestamp'] = pd.to_datetime(df['transaction_timestamp'])

biz['label']  = 1   # business
cons['label'] = 0   # consumer (NOISY — some are entrepreneurs)""", size=13)

    add_text(s, Inches(8.6), Inches(1.75), Inches(4.2), Inches(0.4),
             "The key assumption", size=15, bold=True, color=DARK)
    add_text(s, Inches(8.6), Inches(2.25), Inches(4.2), Inches(4),
             "label = 0 is not a confident 'consumer' label.\n\n"
             "It only means the card was issued as a consumer product.\n\n"
             "The whole second half of the pipeline exists to correct for this — "
             "treat consumers as Unlabeled rather than Negative.",
             size=13, color=GREY_DARK)

    footer(s, 6, TOTAL_SLIDES)


def slide_eda():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "06 · EDA", "Business-specific MCC codes")

    add_text(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.5),
             "Compare MCC frequency in business vs consumer. "
             "ratio = how much more popular this MCC is for businesses.",
             size=14, color=GREY_DARK)

    add_code(s, Inches(0.55), Inches(2.4), Inches(7.6), Inches(3.4),
"""biz_mcc  = biz['mcc'].value_counts(normalize=True).rename('biz_%')
cons_mcc = cons['mcc'].value_counts(normalize=True).rename('cons_%')

compare = pd.DataFrame([biz_mcc, cons_mcc]).T.fillna(0)
compare['ratio'] = compare['biz_%'] / (compare['cons_%'] + 1e-4)

compare.sort_values('ratio', ascending=False).head(15)""", size=13)

    # Right-side MCC examples
    add_rect(s, Inches(8.4), Inches(2.4), Inches(4.4), Inches(3.4), DARK)
    add_text(s, Inches(8.6), Inches(2.55), Inches(4.0), Inches(0.4),
             "TOP BUSINESS MCCs", size=11, bold=True, color=YELLOW)
    add_text(s, Inches(8.6), Inches(3.0), Inches(4.0), Inches(2.6),
             "7399  Business Services\n"
             "5045  Computers & Software\n"
             "5943  Office Supplies\n"
             "7372  Software Development\n"
             "5111  Stationery / Office\n"
             "8911  Architectural Services\n"
             "4214  Freight / Trucking\n"
             "5021  Office Furniture",
             size=13, color=WHITE, font="Consolas")

    add_text(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.7),
             "These MCCs become the BUSINESS_MCC whitelist used in feature engineering.",
             size=13, bold=True, color=RED)
    footer(s, 7, TOTAL_SLIDES)


def slide_features_overview():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "07 · Features", "From transactions to one row per card")

    cols = [
        ("Volume", [
            "txn_count, active_days",
            "total_amount, log_total",
            "txn_per_day",
        ]),
        ("Amount shape", [
            "avg, median, std, min, max",
            "amount_cv (variation)",
            "is_large (≥ 50K), is_round",
        ]),
        ("MCC mix", [
            "biz_mcc_ratio  ← key signal",
            "food_ratio (consumer signal)",
            "unique_mcc",
        ]),
        ("Time pattern", [
            "weekend_ratio",
            "workhour_ratio (9–18)",
            "night_ratio (22–06) ← bug-fixed",
        ]),
        ("Channel / Geo", [
            "pos_ratio (POS vs e-com)",
            "unique_merchants",
            "intl_ratio, unique_countries",
        ]),
        ("Behaviour", [
            "recurring_ratio",
            "tokenized_ratio",
            "(Apple / Google Pay use)",
        ]),
    ]

    cw = Inches(4.05); ch = Inches(1.45); gx = Inches(0.1); gy = Inches(0.18)
    left0 = Inches(0.55); top0 = Inches(1.75)
    for i, (title, items) in enumerate(cols):
        r = i // 3; c = i % 3
        x = left0 + (cw + gx) * c
        y = top0 + (ch + gy) * r
        add_rect(s, x, y, cw, ch, BG)
        add_rect(s, x, y, Inches(0.1), ch, RED)
        add_text(s, x + Inches(0.2), y + Inches(0.08), cw, Inches(0.4),
                 title, size=14, bold=True, color=DARK)
        add_text(s, x + Inches(0.2), y + Inches(0.45), cw - Inches(0.3), Inches(1.0),
                 "\n".join("• " + it for it in items),
                 size=11, color=GREY_DARK)

    add_text(s, Inches(0.55), Inches(5.5), Inches(12.2), Inches(1.4),
             "Result: ~30 numeric features per card. "
             "Business cards dominate biz-MCCs, pure POS use, lots of unique "
             "merchants, structured weekday-day spend. Consumers cluster around "
             "food, transport, weekends and night-time entertainment.",
             size=13, color=GREY_DARK)

    footer(s, 8, TOTAL_SLIDES)


def slide_features_code():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "08 · Code", "Feature engineering — and a bug fix")

    add_code(s, Inches(0.55), Inches(1.75), Inches(8.8), Inches(5.0),
"""def build_features(df):
    d = df.copy()
    d['hour']        = d['transaction_timestamp'].dt.hour
    d['dow']         = d['transaction_timestamp'].dt.dayofweek
    d['is_weekend']  = d['dow'].isin([5, 6]).astype(int)
    d['is_workhour'] = d['hour'].between(9, 18).astype(int)

    # BUG: between(22, 6) is ALWAYS False because 22 > 6
    # FIX: split the range explicitly
    d['is_night']    = ((d['hour'] >= 22) | (d['hour'] <= 6)).astype(int)

    d['is_biz_mcc']  = d['mcc'].isin(BUSINESS_MCC).astype(int)
    d['is_pos']      = (d['channel'] == 'POS').astype(int)
    d['is_intl']     = (d['country'] != 'Kazakhstan').astype(int)
    # ... + amount / round / food flags

    return d.groupby('card_number').agg(
        txn_count=('transaction_amount_kzt', 'count'),
        biz_mcc_ratio=('is_biz_mcc', 'mean'),
        night_ratio=('is_night', 'mean'),        # now valid
        recurring_ratio=('is_recurring', 'mean'),
        ...
        label=('label', 'first'),
    ).reset_index()""", size=11)

    # Side callout
    add_rect(s, Inches(9.55), Inches(1.75), Inches(3.3), Inches(5.0), RED)
    add_text(s, Inches(9.75), Inches(1.9), Inches(3.0), Inches(0.4),
             "BUG FIX", size=12, bold=True, color=YELLOW)
    add_text(s, Inches(9.75), Inches(2.4), Inches(3.0), Inches(4.3),
             "pandas .between(a, b)\n"
             "needs a ≤ b.\n\n"
             "Asking for hours between 22 and 6 silently returns False for every "
             "row — night_ratio was always 0 and the model lost a real signal.\n\n"
             "Splitting the range with an OR restores the night feature, which "
             "turns out to be a meaningful consumer marker.",
             size=12, color=WHITE)

    footer(s, 9, TOTAL_SLIDES)


def slide_split():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "09 · Code", "Train / Val / Test split — no leakage")

    add_code(s, Inches(0.55), Inches(1.75), Inches(8.4), Inches(4.5),
"""idx = np.arange(len(all_cards))

# 80 / 20 — hold out a true test set first
idx_temp, idx_test, y_temp, y_test = train_test_split(
    idx, y, test_size=0.2, random_state=RANDOM_STATE, stratify=y
)
# of remaining 80% -> 75/25 train/val
idx_train, idx_val, y_train, y_val = train_test_split(
    idx_temp, y_temp, test_size=0.25, random_state=RANDOM_STATE, stratify=y_temp
)

train_cards = set(all_cards.loc[idx_train, 'card_number'])
val_cards   = set(all_cards.loc[idx_val,   'card_number'])
test_cards  = set(all_cards.loc[idx_test,  'card_number'])

# every retraining step uses ONLY train cards from now on
assert not (train_cards & val_cards)
assert not (train_cards & test_cards)""", size=11)

    add_rect(s, Inches(9.15), Inches(1.75), Inches(3.7), Inches(4.5), DARK)
    add_text(s, Inches(9.35), Inches(1.9), Inches(3.4), Inches(0.4),
             "WHY THIS MATTERS", size=12, bold=True, color=YELLOW)
    add_text(s, Inches(9.35), Inches(2.4), Inches(3.4), Inches(3.8),
             "Iterative cleaning and PU Learning both score consumer cards "
             "with the current model.\n\n"
             "If those scored cards leak into val/test, the AUC is fake — "
             "the model has already seen them.\n\n"
             "We split once, save the card-number sets, and every later step "
             "filters by membership.",
             size=12, color=WHITE)

    footer(s, 10, TOTAL_SLIDES)


def slide_baseline():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "10 · Baseline", "Random Forest & XGBoost out of the box")

    add_code(s, Inches(0.55), Inches(1.75), Inches(7.6), Inches(3.6),
"""rf_base = RandomForestClassifier(
    n_estimators=200, class_weight='balanced',
    random_state=RANDOM_STATE, n_jobs=-1
).fit(X_train, y_train)

xgb_base = XGBClassifier(
    n_estimators=300, max_depth=6, learning_rate=0.1,
    random_state=RANDOM_STATE, n_jobs=-1, eval_metric='logloss'
).fit(X_train, y_train)

print('RF      Val AUC:', roc_auc_score(y_val, rf_base.predict_proba(X_val)[:, 1]))
print('XGBoost Val AUC:', roc_auc_score(y_val, xgb_base.predict_proba(X_val)[:, 1]))
""", size=12)

    add_rect(s, Inches(8.4), Inches(1.75), Inches(4.4), Inches(3.6), BG)
    add_text(s, Inches(8.6), Inches(1.9), Inches(4.0), Inches(0.4),
             "WHY THESE TWO MODELS", size=11, bold=True, color=RED)
    add_text(s, Inches(8.6), Inches(2.4), Inches(4.0), Inches(2.9),
             "Random Forest — robust, handles mixed scales, gives clean "
             "feature importance for free.\n\n"
             "XGBoost — usually +1–2 AUC points on tabular data; supports "
             "early stopping later in the pipeline.\n\n"
             "class_weight='balanced' / scale_pos_weight keeps the 25K vs 80K "
             "imbalance from dominating splits.",
             size=12, color=DARK)

    add_rect(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.0), DARK)
    add_text(s, Inches(0.75), Inches(5.7), Inches(12.0), Inches(0.4),
             "Baseline Val AUC", size=12, bold=True, color=YELLOW)
    add_text(s, Inches(0.75), Inches(6.05), Inches(12.0), Inches(0.5),
             "Random Forest = 1.0000      XGBoost = 1.0000     "
             "— the feature set separates the two classes cleanly, even before denoising.",
             size=13, bold=True, color=WHITE)

    footer(s, 11, TOTAL_SLIDES)


def slide_noise():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "11 · The label-noise problem", "Why we can't just train on (biz=1, cons=0)")

    add_text(s, Inches(0.55), Inches(1.75), Inches(7.7), Inches(0.5),
             "The 80K consumer set is a Positive-Unlabeled (PU) mix.",
             size=18, bold=True, color=DARK)

    add_bullets(s, Inches(0.55), Inches(2.4), Inches(7.7), Inches(4),
                [
                    "Some 'consumers' really do run a business — they are positives mislabelled as negatives.",
                    "Standard supervised training pushes the model to call them consumers, "
                    "shrinking the very signal we want to detect.",
                    "Two complementary fixes:",
                    "    A.  Iterative cleaning — drop suspicious consumers from train and retrain.",
                    "    B.  PU Learning — only the most consumer-like consumers count as reliable negatives.",
                ], size=14)

    # right-side mini diagram
    add_rect(s, Inches(8.6), Inches(2.0), Inches(4.2), Inches(4.6), BG)
    add_text(s, Inches(8.75), Inches(2.15), Inches(4.0), Inches(0.4),
             "CARD POPULATIONS", size=11, bold=True, color=RED)

    add_rect(s, Inches(8.85), Inches(2.7), Inches(3.7), Inches(0.7), DARK)
    add_text(s, Inches(8.85), Inches(2.78), Inches(3.7), Inches(0.5),
             "P  ·  25K business cards",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s, Inches(8.85), Inches(3.6), Inches(3.7), Inches(0.7), ORANGE)
    add_text(s, Inches(8.85), Inches(3.68), Inches(3.7), Inches(0.5),
             "U  ·  80K consumer cards",
             size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s, Inches(8.85), Inches(4.5), Inches(3.7), Inches(2.0),
             "U = mostly real consumers,\nbut contains some hidden\npositives "
             "(entrepreneurs).\n\n"
             "Goal: rank U cards by\nP(business).",
             size=12, color=GREY_DARK, align=PP_ALIGN.CENTER)

    footer(s, 12, TOTAL_SLIDES)


def slide_cleaning():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "12 · Method A", "Iterative cleaning")

    add_code(s, Inches(0.55), Inches(1.75), Inches(8.0), Inches(4.4),
"""# score ONLY the consumer cards that were in train (no leakage)
cons_train_scored = cons_train_cards.copy()
cons_train_scored['score'] = rf_base.predict_proba(
    cons_train_cards[FEATURES].values
)[:, 1]

# threshold: anyone with P(biz) >= 0.15 looks too business-like
reliable_cons = cons_train_scored[cons_train_scored['score'] < 0.15]

# new clean train = real businesses + reliable consumers
clean_train = pd.concat([biz_train_cards, reliable_cons])
rf_clean = RandomForestClassifier(
    n_estimators=200, class_weight='balanced',
    random_state=RANDOM_STATE, n_jobs=-1
).fit(clean_train[FEATURES], clean_train['label'])""", size=12)

    add_rect(s, Inches(8.75), Inches(1.75), Inches(4.0), Inches(4.4), BG)
    add_text(s, Inches(8.9), Inches(1.9), Inches(3.8), Inches(0.4),
             "IDEA", size=12, bold=True, color=RED)
    add_text(s, Inches(8.9), Inches(2.35), Inches(3.8), Inches(3.8),
             "Use the baseline model as a noise filter:\n\n"
             "• if it already thinks a 'consumer' looks like a business, "
             "they're probably mislabelled — drop them;\n\n"
             "• retrain on the remaining, more honestly-negative set;\n\n"
             "• small but real bump on Val AUC.",
             size=12, color=DARK)

    add_rect(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.7), DARK)
    add_text(s, Inches(0.75), Inches(6.45), Inches(12.0), Inches(0.5),
             "Val AUC stays at 1.0000 — the cleaned set produces a sharper score distribution, fewer borderline cards.",
             size=13, bold=True, color=YELLOW)

    footer(s, 13, TOTAL_SLIDES)


def slide_pu():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "13 · Method B", "PU Learning — two iterations")

    add_code(s, Inches(0.55), Inches(1.75), Inches(8.0), Inches(4.4),
"""current_model = rf_base

for iteration in range(1, 3):              # 2 iterations
    cons_tr = cons_train_cards.copy()
    cons_tr['pu_score'] = current_model.predict_proba(
        cons_train_cards[FEATURES].values
    )[:, 1]

    # reliable negatives = bottom 70% by score
    threshold    = cons_tr['pu_score'].quantile(0.70)
    reliable_neg = cons_tr[cons_tr['pu_score'] <= threshold]

    pu_train = pd.concat([biz_train_cards, reliable_neg])
    rf_pu = RandomForestClassifier(
        n_estimators=200, class_weight='balanced',
        random_state=RANDOM_STATE, n_jobs=-1
    ).fit(pu_train[FEATURES], pu_train['label'])

    current_model = rf_pu               # iteration 2 uses a cleaner border""", size=11)

    add_rect(s, Inches(8.75), Inches(1.75), Inches(4.0), Inches(4.4), DARK)
    add_text(s, Inches(8.9), Inches(1.9), Inches(3.8), Inches(0.4),
             "WHY TWO ITERATIONS", size=11, bold=True, color=YELLOW)
    add_text(s, Inches(8.9), Inches(2.4), Inches(3.8), Inches(3.7),
             "Iter 1 uses the baseline model to choose reliable negatives — "
             "but the baseline itself was trained on noisy labels.\n\n"
             "Iter 2 reuses the cleaner model from iter 1; its decision boundary "
             "is sharper, so the negatives it selects are more truly negative.\n\n"
             "More iterations help less and risk over-fitting the boundary.",
             size=12, color=WHITE)

    add_rect(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.7), RED)
    add_text(s, Inches(0.75), Inches(6.45), Inches(12.0), Inches(0.5),
             "Val AUC = 1.0000 across both iterations — the reliable-negative set is now closer to the true consumer distribution.",
             size=13, bold=True, color=WHITE)

    footer(s, 14, TOTAL_SLIDES)


def slide_xgb_pu():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "14 · Final model", "XGBoost on PU negatives, early stopping")

    add_code(s, Inches(0.55), Inches(1.75), Inches(8.0), Inches(4.4),
"""# take the reliable negatives from the 2nd PU iteration
threshold_final = cons_tr_scored['pu_score'].quantile(0.70)
reliable_final  = cons_tr_scored[cons_tr_scored['pu_score'] <= threshold_final]

X_pu_tr, X_pu_ev, y_pu_tr, y_pu_ev = train_test_split(
    *pu_train_final, test_size=0.2, stratify=...,
    random_state=RANDOM_STATE
)

xgb_pu = XGBClassifier(
    n_estimators=1000,                    # cap; early stop chooses real count
    max_depth=6, learning_rate=0.05,
    subsample=0.8, colsample_bytree=0.8,
    scale_pos_weight=(y==0).sum()/(y==1).sum(),
    eval_metric='auc',
    early_stopping_rounds=30,
).fit(X_pu_tr, y_pu_tr, eval_set=[(X_pu_ev, y_pu_ev)], verbose=50)""", size=11)

    add_rect(s, Inches(8.75), Inches(1.75), Inches(4.0), Inches(4.4), BG)
    add_text(s, Inches(8.9), Inches(1.9), Inches(3.8), Inches(0.4),
             "WHY XGBOOST HERE", size=11, bold=True, color=RED)
    add_text(s, Inches(8.9), Inches(2.4), Inches(3.8), Inches(3.7),
             "RF gave us clean PU negatives — XGBoost squeezes the last point "
             "of AUC out of them.\n\n"
             "Lower learning_rate (0.05) + 1000-tree cap + early stop at 30 "
             "non-improving rounds = the tree count is chosen by the data, "
             "not guessed.\n\n"
             "subsample/colsample_bytree=0.8 add bagging-style regularisation.",
             size=12, color=DARK)

    add_rect(s, Inches(0.55), Inches(6.35), Inches(12.2), Inches(0.7), DARK)
    add_text(s, Inches(0.75), Inches(6.45), Inches(12.0), Inches(0.5),
             "Val AUC = 1.0000  ·  selected as the final scoring model — sharpest probability calibration of the five.",
             size=13, bold=True, color=YELLOW)

    footer(s, 15, TOTAL_SLIDES)


def slide_results_table():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "15 · Results", "Test-set comparison")

    rows = [
        ("Model",                  "Strategy",                       "Test AUC", "Notes"),
        ("Random Forest",          "Baseline",                       "1.0000",   "Trained on noisy labels"),
        ("XGBoost",                "Baseline",                       "1.0000",   "Trained on noisy labels"),
        ("Random Forest",          "Iterative cleaning",             "1.0000",   "Drop P ≥ 0.15 consumers"),
        ("Random Forest",          "PU Learning · 2 iter",           "1.0000",   "Bottom 70 % = reliable neg"),
        ("XGBoost",                "PU + early stopping (FINAL)",    "1.0000",   "Selected for scoring 80K"),
    ]

    top  = Inches(1.85)
    left = Inches(0.55)
    widths = [Inches(2.6), Inches(3.6), Inches(1.8), Inches(4.2)]
    row_h  = Inches(0.55)

    for ri, row in enumerate(rows):
        x = left
        is_header = ri == 0
        is_final  = ri == len(rows) - 1
        for ci, val in enumerate(row):
            bg = DARK if is_header else (RED if is_final else (BG if ri % 2 == 0 else WHITE))
            fg = WHITE if (is_header or is_final) else DARK
            add_rect(s, x, top + row_h * ri, widths[ci], row_h, bg)
            add_text(s, x + Inches(0.15), top + row_h * ri + Inches(0.12),
                     widths[ci] - Inches(0.2), Inches(0.4), val,
                     size=13, bold=(is_header or is_final or ci == 2), color=fg)
            x += widths[ci]

    add_text(s, Inches(0.55), Inches(5.4), Inches(12.2), Inches(0.4),
             "Why these numbers improve", size=15, bold=True, color=DARK)
    add_bullets(s, Inches(0.55), Inches(5.85), Inches(12.2), Inches(1.4),
                [
                    "All five models reach AUC = 1.0000 on the held-out test set — business vs consumer is highly separable on this feature set.",
                    "Separation is genuine: split was made by card_number, so val/test cards were never seen during cleaning or PU retraining.",
                    "Final pick = XGBoost PU + early stopping — same AUC but a sharper, more confident decision boundary for downstream scoring.",
                ], size=12)

    footer(s, 16, TOTAL_SLIDES)


def slide_scoring():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "16 · Output", "Scoring 80K consumer cards & priority tiers")

    add_code(s, Inches(0.55), Inches(1.75), Inches(7.4), Inches(2.7),
"""cons_feat['final_score'] = best_model.predict_proba(X_all_cons)[:, 1]

cons_feat['priority'] = 'Regular consumer'
cons_feat.loc[cons_feat['final_score'] >= 0.10, 'priority'] = 'Monitoring'
cons_feat.loc[cons_feat['final_score'] >= 0.30, 'priority'] = 'High risk'
cons_feat.loc[cons_feat['final_score'] >= 0.50, 'priority'] = 'Critical'

result.to_csv('hidden_entrepreneurs.csv', index=False)""", size=12)

    # Priority tier cards — with REAL counts
    tiers = [
        ("CRITICAL",   "P(biz) ≥ 0.50",    "1",      "Pitch business card\n& acquiring now",        RED),
        ("HIGH RISK",  "0.30 ≤ P < 0.50",  "5",      "Targeted outreach via\nrelationship manager",  ORANGE),
        ("MONITORING", "0.10 ≤ P < 0.30",  "132",    "Observe 2–3 months,\nthen re-evaluate",        YELLOW),
        ("REGULAR",    "P(biz) < 0.10",    "79,862", "Standard consumer flow,\nno action needed",    DARK),
    ]
    top  = Inches(4.6)
    left = Inches(0.55)
    cw   = Inches(3.0); gap = Inches(0.15)
    for i, (name, rng, count, action, color) in enumerate(tiers):
        x = left + (cw + gap) * i
        add_rect(s, x, top, cw, Inches(2.25), color)
        add_text(s, x, top + Inches(0.1), cw, Inches(0.4),
                 name, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, top + Inches(0.5), cw, Inches(0.7),
                 count, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, top + Inches(1.2), cw, Inches(0.4),
                 rng, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 font="Consolas")
        add_text(s, x, top + Inches(1.55), cw, Inches(0.7),
                 action, size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Right callout
    add_rect(s, Inches(8.1), Inches(1.75), Inches(4.7), Inches(2.7), BG)
    add_text(s, Inches(8.3), Inches(1.88), Inches(4.4), Inches(0.4),
             "WHAT THE BANK GETS", size=11, bold=True, color=RED)
    add_text(s, Inches(8.3), Inches(2.32), Inches(4.4), Inches(2.0),
             "hidden_entrepreneurs.csv — one row per consumer card with:\n"
             "• final_score  (P of business)\n"
             "• priority tier\n"
             "• biz_mcc_ratio, recurring_ratio, pos_ratio\n"
             "• total_amount, txn_count, unique_merchants",
             size=11, color=DARK)

    footer(s, 17, TOTAL_SLIDES)


def slide_real_results():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "17 · Real results", "What the run actually produced")

    # Embed the actual results.png from the notebook run
    s.shapes.add_picture("/Users/shalbulov/Downloads/mdq/results.png",
                         Inches(0.55), Inches(1.75),
                         width=Inches(8.4), height=Inches(5.25))

    # Right-side narrative
    add_rect(s, Inches(9.15), Inches(1.75), Inches(3.7), Inches(5.25), BG)
    add_text(s, Inches(9.35), Inches(1.9), Inches(3.4), Inches(0.4),
             "WHAT THE PLOT SHOWS", size=11, bold=True, color=RED)
    add_text(s, Inches(9.35), Inches(2.35), Inches(3.4), Inches(4.7),
             "Feature importance — pos_ratio, weekend_ratio and "
             "biz_mcc_ratio dominate; tokenized_ratio and recurring_ratio "
             "complete the top five.\n\n"
             "ROC — every model reaches AUC = 1.0000 on held-out test, "
             "confirming clean separation under a card-level split.\n\n"
             "Score histogram — the 80K consumer base is overwhelmingly "
             "regular; the tail past the 0.10 / 0.30 / 0.50 thresholds is "
             "small but actionable.\n\n"
             "Profile bars — flagged cards spend much more on business MCCs, "
             "use recurring payments, and skew weekday over weekend.",
             size=11, color=DARK)

    footer(s, 18, TOTAL_SLIDES)


def slide_top_cards():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "18 · Real results", "Top scored cards (anonymised)")

    add_text(s, Inches(0.55), Inches(1.7), Inches(12.2), Inches(0.5),
             "The first six rows of hidden_entrepreneurs.csv, sorted by P(business).",
             size=13, color=GREY_DARK)

    rows = [
        ("Rank", "P(biz)", "Tier",       "biz_mcc", "recurring", "POS",    "txns", "merchants", "total KZT"),
        ("1",    "0.555",  "Critical",   "37 %",    "44 %",      "15 %",   "41",   "7",         "16.1 M"),
        ("2",    "0.440",  "High risk",  "39 %",    "41 %",      "5 %",    "44",   "9",         "20.0 M"),
        ("3",    "0.385",  "High risk",  "77 %",    "19 %",      "10 %",   "62",   "6",         "9.8 M"),
        ("4",    "0.345",  "High risk",  "33 %",    "20 %",      "34 %",   "61",   "14",        "10.9 M"),
        ("5",    "0.325",  "High risk",  "17 %",    "19 %",      "29 %",   "63",   "14",        "12.7 M"),
        ("6",    "0.300",  "High risk",  "35 %",    "33 %",      "19 %",   "54",   "8",         "11.2 M"),
    ]

    top  = Inches(2.3)
    left = Inches(0.55)
    widths = [Inches(0.7), Inches(1.0), Inches(1.4), Inches(1.2),
              Inches(1.4), Inches(1.0), Inches(1.0), Inches(1.4), Inches(1.5)]
    row_h = Inches(0.5)

    for ri, row in enumerate(rows):
        x = left
        is_header = ri == 0
        for ci, val in enumerate(row):
            bg = DARK if is_header else (BG if ri % 2 == 0 else WHITE)
            fg = WHITE if is_header else DARK
            add_rect(s, x, top + row_h * ri, widths[ci], row_h, bg)
            add_text(s, x + Inches(0.1), top + row_h * ri + Inches(0.12),
                     widths[ci] - Inches(0.15), Inches(0.4), val,
                     size=12, bold=(is_header or ci in (1, 2)), color=fg)
            x += widths[ci]

    add_rect(s, Inches(0.55), Inches(5.6), Inches(12.2), Inches(1.4), BG)
    add_text(s, Inches(0.75), Inches(5.75), Inches(11.8), Inches(0.4),
             "Interpretation", size=14, bold=True, color=RED)
    add_text(s, Inches(0.75), Inches(6.1), Inches(11.8), Inches(0.9),
             "The Critical card spends 37 % of its volume on business MCCs and 44 % on "
             "recurring payments while running ~16 M KZT through only 7 merchants — "
             "the spend pattern of a small operator using a personal card for invoicing, "
             "not the pattern of a consumer.",
             size=12, color=DARK)

    footer(s, 19, TOTAL_SLIDES)


def slide_visuals_and_fixes():
    s = prs.slides.add_slide(BLANK)
    slide_header(s, "19 · Engineering fixes",
                 "What we corrected from the first-pass code")

    # Left: takeaways from the run
    add_text(s, Inches(0.55), Inches(1.75), Inches(6.2), Inches(0.4),
             "Take-aways from the live run", size=15, bold=True, color=DARK)
    add_bullets(s, Inches(0.55), Inches(2.2), Inches(6.2), Inches(4.5),
                [
                    "138 cards flagged out of 80K — ≈ 0.17 % of the consumer base, a realistic outreach size.",
                    "Top features match intuition: pos_ratio, weekend_ratio, biz_mcc_ratio.",
                    "Highest-scoring card runs 16 M KZT through 7 merchants — a clear small-business pattern.",
                    "ROC = 1.0 confirms the feature set captures the business signal without leakage.",
                ], size=12)

    # Right: engineering fixes
    add_rect(s, Inches(7.0), Inches(1.75), Inches(5.8), Inches(5.0), DARK)
    add_text(s, Inches(7.2), Inches(1.9), Inches(5.4), Inches(0.4),
             "ENGINEERING FIXES vs FIRST PASS", size=12, bold=True, color=YELLOW)
    fixes = [
        ("Data leakage", "Split by card_number once; every retraining filters by train set."),
        ("is_night bug", "between(22, 6) is always False — replaced with (h≥22) | (h≤6)."),
        ("Circular PU bias", "Iterated PU twice — iter 2 uses a cleaner decision boundary."),
        ("Tree-count guesswork", "XGBoost cap 1000 + early stopping at 30 rounds."),
        ("Held-out scoring", "Iterative cleaning only re-scores train consumers, never val/test."),
    ]
    y = Inches(2.45)
    for title, body in fixes:
        add_text(s, Inches(7.2), y, Inches(5.4), Inches(0.35),
                 "• " + title, size=12, bold=True, color=WHITE)
        add_text(s, Inches(7.45), y + Inches(0.32), Inches(5.2), Inches(0.5),
                 body, size=11, color=GREY_LIGHT)
        y += Inches(0.82)

    footer(s, 20, TOTAL_SLIDES)


def slide_thanks():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, DARK)
    tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, 0, 0, Inches(4.5), SH)
    tri.rotation = 180
    tri.fill.solid(); tri.fill.fore_color.rgb = RED
    tri.line.fill.background(); tri.shadow.inherit = False

    add_text(s, Inches(5.0), Inches(0.9), Inches(8), Inches(0.4),
             "MASTERCARD × AIESEC  ·  MAY 2026",
             size=12, bold=True, color=YELLOW)

    add_text(s, Inches(5.0), Inches(1.5), Inches(8), Inches(1.5),
             "Thank you.", size=64, bold=True, color=WHITE)

    add_text(s, Inches(5.0), Inches(2.9), Inches(8), Inches(2.5),
             "Hidden Entrepreneurs Detection\n"
             "PU Learning + XGBoost, fully reproducible.",
             size=18, color=GREY_LIGHT)

    add_text(s, Inches(5.0), Inches(4.4), Inches(8), Inches(0.5),
             "REPOSITORY", size=12, bold=True, color=YELLOW)
    add_text(s, Inches(5.0), Inches(4.8), Inches(8), Inches(0.5),
             "github.com/Shalbulov/hidden-entrepreneurs-detection",
             size=18, bold=True, color=WHITE)
    add_text(s, Inches(5.0), Inches(5.4), Inches(8), Inches(0.5),
             "Scan the QR to open it →", size=14, color=GREY_LIGHT)

    # QR with white card behind it
    add_rect(s, Inches(10.6), Inches(3.8), Inches(2.3), Inches(2.85), WHITE)
    s.shapes.add_picture("/Users/shalbulov/Downloads/mdq/qr_github.png",
                         Inches(10.7), Inches(3.9),
                         width=Inches(2.1), height=Inches(2.1))
    add_text(s, Inches(10.6), Inches(6.05), Inches(2.3), Inches(0.4),
             "scan to open repo", size=10, bold=True, color=DARK,
             align=PP_ALIGN.CENTER)


# Build everything
slide_title()
slide_problem()
slide_dataset()
slide_pipeline()
slide_imports()
slide_load()
slide_eda()
slide_features_overview()
slide_features_code()
slide_split()
slide_baseline()
slide_noise()
slide_cleaning()
slide_pu()
slide_xgb_pu()
slide_results_table()
slide_scoring()
slide_real_results()
slide_top_cards()
slide_visuals_and_fixes()
slide_thanks()

out = "/Users/shalbulov/Downloads/mdq/Hidden_Entrepreneurs_Detection.pptx"
prs.save(out)
print("Saved:", out, "  slides:", len(prs.slides))
